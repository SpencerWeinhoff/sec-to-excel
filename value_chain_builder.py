"""Build industry value chain PowerPoint presentations.

Ported from value-chain-generator.js layout engine to Python/python-pptx.
Uses pre-built value chain data (no external API calls).
"""

import os
import tempfile
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================================
# COLOR PALETTES — keyed by industry keyword
# ============================================================================

PALETTES = {
    "default": {
        "titleBg": "1B2A4A",
        "titleText": "FFFFFF",
        "slideBg": "F5F6FA",
        "cardBg": "FFFFFF",
        "cardBorder": "E2E5EC",
        "accent": "3B6FE0",
        "accentLight": "EBF0FB",
        "stageText": "1B2A4A",
        "bodyText": "4A5568",
        "arrowColor": "3B6FE0",
        "badgeBg": "EBF0FB",
        "badgeText": "2D5BC4",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "ai": {
        "titleBg": "0F0F2D",
        "titleText": "FFFFFF",
        "slideBg": "F4F3FA",
        "cardBg": "FFFFFF",
        "cardBorder": "DDD8F0",
        "accent": "7C3AED",
        "accentLight": "EDE9FE",
        "stageText": "1E1B4B",
        "bodyText": "4A4568",
        "arrowColor": "7C3AED",
        "badgeBg": "EDE9FE",
        "badgeText": "5B21B6",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "semiconductor": {
        "titleBg": "111827",
        "titleText": "FFFFFF",
        "slideBg": "F3F4F6",
        "cardBg": "FFFFFF",
        "cardBorder": "D1D5DB",
        "accent": "2563EB",
        "accentLight": "DBEAFE",
        "stageText": "111827",
        "bodyText": "4B5563",
        "arrowColor": "2563EB",
        "badgeBg": "DBEAFE",
        "badgeText": "1D4ED8",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "energy": {
        "titleBg": "14532D",
        "titleText": "FFFFFF",
        "slideBg": "F0FDF4",
        "cardBg": "FFFFFF",
        "cardBorder": "BBF7D0",
        "accent": "16A34A",
        "accentLight": "DCFCE7",
        "stageText": "14532D",
        "bodyText": "3F6B52",
        "arrowColor": "16A34A",
        "badgeBg": "DCFCE7",
        "badgeText": "15803D",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "oilgas": {
        "titleBg": "1C1917",
        "titleText": "FFFFFF",
        "slideBg": "FAF9F6",
        "cardBg": "FFFFFF",
        "cardBorder": "D6D3D1",
        "accent": "B45309",
        "accentLight": "FEF3C7",
        "stageText": "1C1917",
        "bodyText": "57534E",
        "arrowColor": "B45309",
        "badgeBg": "FEF3C7",
        "badgeText": "92400E",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "pharma": {
        "titleBg": "1E1B4B",
        "titleText": "FFFFFF",
        "slideBg": "FAF5FF",
        "cardBg": "FFFFFF",
        "cardBorder": "E9D5FF",
        "accent": "9333EA",
        "accentLight": "F3E8FF",
        "stageText": "1E1B4B",
        "bodyText": "5B4A6E",
        "arrowColor": "9333EA",
        "badgeBg": "F3E8FF",
        "badgeText": "7E22CE",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "fintech": {
        "titleBg": "0C4A6E",
        "titleText": "FFFFFF",
        "slideBg": "F0F9FF",
        "cardBg": "FFFFFF",
        "cardBorder": "BAE6FD",
        "accent": "0284C7",
        "accentLight": "E0F2FE",
        "stageText": "0C4A6E",
        "bodyText": "475569",
        "arrowColor": "0284C7",
        "badgeBg": "E0F2FE",
        "badgeText": "0369A1",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "cyber": {
        "titleBg": "18181B",
        "titleText": "FFFFFF",
        "slideBg": "F4F4F5",
        "cardBg": "FFFFFF",
        "cardBorder": "D4D4D8",
        "accent": "DC2626",
        "accentLight": "FEE2E2",
        "stageText": "18181B",
        "bodyText": "52525B",
        "arrowColor": "DC2626",
        "badgeBg": "FEE2E2",
        "badgeText": "B91C1C",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "aerospace": {
        "titleBg": "0F172A",
        "titleText": "FFFFFF",
        "slideBg": "F1F5F9",
        "cardBg": "FFFFFF",
        "cardBorder": "CBD5E1",
        "accent": "334155",
        "accentLight": "E2E8F0",
        "stageText": "0F172A",
        "bodyText": "475569",
        "arrowColor": "334155",
        "badgeBg": "E2E8F0",
        "badgeText": "1E293B",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "ecommerce": {
        "titleBg": "431407",
        "titleText": "FFFFFF",
        "slideBg": "FFF7ED",
        "cardBg": "FFFFFF",
        "cardBorder": "FED7AA",
        "accent": "EA580C",
        "accentLight": "FFEDD5",
        "stageText": "431407",
        "bodyText": "6B5344",
        "arrowColor": "EA580C",
        "badgeBg": "FFEDD5",
        "badgeText": "C2410C",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "cloud": {
        "titleBg": "172554",
        "titleText": "FFFFFF",
        "slideBg": "EFF6FF",
        "cardBg": "FFFFFF",
        "cardBorder": "BFDBFE",
        "accent": "2563EB",
        "accentLight": "DBEAFE",
        "stageText": "172554",
        "bodyText": "3B5998",
        "arrowColor": "2563EB",
        "badgeBg": "DBEAFE",
        "badgeText": "1D4ED8",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "ev": {
        "titleBg": "052E16",
        "titleText": "FFFFFF",
        "slideBg": "F0FDF4",
        "cardBg": "FFFFFF",
        "cardBorder": "BBF7D0",
        "accent": "059669",
        "accentLight": "D1FAE5",
        "stageText": "052E16",
        "bodyText": "3F6B52",
        "arrowColor": "059669",
        "badgeBg": "D1FAE5",
        "badgeText": "047857",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
    "telecom": {
        "titleBg": "1E1B4B",
        "titleText": "FFFFFF",
        "slideBg": "EEF2FF",
        "cardBg": "FFFFFF",
        "cardBorder": "C7D2FE",
        "accent": "4F46E5",
        "accentLight": "E0E7FF",
        "stageText": "1E1B4B",
        "bodyText": "4338CA",
        "arrowColor": "4F46E5",
        "badgeBg": "E0E7FF",
        "badgeText": "3730A3",
        "headerFont": "Georgia",
        "bodyFont": "Calibri",
    },
}


def _get_palette(industry_str):
    """Match industry string to a color palette."""
    key = industry_str.lower()
    for k, v in PALETTES.items():
        if k != "default" and k in key:
            return v
    # Fuzzy matching
    if any(w in key for w in ("renew", "solar", "wind", "clean")):
        return PALETTES["energy"]
    if any(w in key for w in ("oil", "gas", "petro")):
        return PALETTES["oilgas"]
    if any(w in key for w in ("chip", "semi", "wafer")):
        return PALETTES["semiconductor"]
    if any(w in key for w in ("drug", "bio", "pharma", "therapeut")):
        return PALETTES["pharma"]
    if any(w in key for w in ("artificial", "machine learn", " ai", "llm", "neural")):
        return PALETTES["ai"]
    if any(w in key for w in ("bank", "fintech", "payment", "lending")):
        return PALETTES["fintech"]
    if any(w in key for w in ("cyber", "security", "threat")):
        return PALETTES["cyber"]
    if any(w in key for w in ("aero", "space", "defense", "satellite")):
        return PALETTES["aerospace"]
    if any(w in key for w in ("ecommerce", "e-commerce", "retail", "shop")):
        return PALETTES["ecommerce"]
    if any(w in key for w in ("cloud", "saas", "devops")):
        return PALETTES["cloud"]
    if any(w in key for w in ("electric vehicle", " ev ", "battery", "lithium")):
        return PALETTES["ev"]
    if any(w in key for w in ("telecom", "5g", "network", "connect")):
        return PALETTES["telecom"]
    return PALETTES["default"]


def _hex_to_rgb(hex_str):
    """Convert hex string to RGBColor."""
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


# ============================================================================
# LAYOUT ENGINE — all coordinates computed mathematically
# ============================================================================

SLIDE_W = 10.0      # inches
SLIDE_H = 5.625     # inches (16:9)
MARGIN = 0.5        # edge margin


class OverviewLayout:
    """Compute positions for the overview (hero flow) slide."""
    TITLE_Y = 0.2
    TITLE_H = 0.6
    TITLE_FONT_SIZE = 26
    FLOW_AREA_TOP = 1.0
    FLOW_AREA_BOTTOM = 4.8
    CARD_PADDING_X = 0.55
    ARROW_WIDTH = 0.22
    ARROW_GAP = 0.08
    ICON_SIZE = 0.4
    ICON_TOP_PAD = 0.25
    NAME_TOP_PAD = 0.08
    NAME_FONT_SIZE = 10
    DESC_TOP_PAD = 0.1
    DESC_FONT_SIZE = 8
    STAGE_NUMBER_SIZE = 8

    @classmethod
    def get_card_layout(cls, stage_count):
        available_w = SLIDE_W - 2 * cls.CARD_PADDING_X
        total_arrow_space = (stage_count - 1) * (cls.ARROW_WIDTH + 2 * cls.ARROW_GAP)
        total_card_w = available_w - total_arrow_space
        card_w = total_card_w / stage_count
        card_h = cls.FLOW_AREA_BOTTOM - cls.FLOW_AREA_TOP
        step = card_w + cls.ARROW_WIDTH + 2 * cls.ARROW_GAP
        return card_w, card_h, step, cls.CARD_PADDING_X

    @classmethod
    def get_card_position(cls, index, stage_count):
        card_w, card_h, step, start_x = cls.get_card_layout(stage_count)
        x = start_x + index * step
        y = cls.FLOW_AREA_TOP
        return x, y, card_w, card_h

    @classmethod
    def get_arrow_position(cls, index, stage_count):
        card_w, _, step, start_x = cls.get_card_layout(stage_count)
        card_end_x = start_x + index * step + card_w + cls.ARROW_GAP
        mid_y = cls.FLOW_AREA_TOP + (cls.FLOW_AREA_BOTTOM - cls.FLOW_AREA_TOP) / 2
        return card_end_x, mid_y - 0.1, cls.ARROW_WIDTH, 0.2


class DeepDiveLayout:
    """Compute positions for deep-dive (per-stage) slides."""
    STAGE_NUMBER_Y = 0.3
    STAGE_NUMBER_SIZE = 12
    TITLE_Y = 0.35
    TITLE_H = 0.5
    TITLE_FONT_SIZE = 28
    TITLE_X = 1.2
    DESC_Y = 0.95
    DESC_H = 0.55
    DESC_FONT_SIZE = 14
    ACCENT_BAR_Y = 1.65
    ACCENT_BAR_H = 0.04
    ACCENT_BAR_W = 9.0
    PLAYERS_LABEL_Y = 1.9
    PLAYERS_LABEL_FONT_SIZE = 14
    PLAYER_GRID_TOP = 2.3
    PLAYER_CARD_H = 1.1
    PLAYER_CARD_GAP = 0.25

    @classmethod
    def get_player_grid(cls, player_count):
        cols = min(player_count, 4)
        grid_w = SLIDE_W - 2 * MARGIN
        card_w = (grid_w - (cols - 1) * cls.PLAYER_CARD_GAP) / cols
        return cols, card_w

    @classmethod
    def get_player_card_position(cls, index, player_count):
        cols, card_w = cls.get_player_grid(player_count)
        col = index % cols
        row = index // cols
        x = MARGIN + col * (card_w + cls.PLAYER_CARD_GAP)
        y = cls.PLAYER_GRID_TOP + row * (cls.PLAYER_CARD_H + cls.PLAYER_CARD_GAP)
        return x, y, card_w, cls.PLAYER_CARD_H


class SummaryLayout:
    TITLE_Y = 0.3
    TITLE_H = 0.6
    TITLE_FONT_SIZE = 26
    TABLE_TOP = 1.2
    TABLE_X = 0.5
    TABLE_W = 9.0
    ROW_H = 0.55
    COMPANY_COL_W = 2.5


# ============================================================================
# SLIDE BUILDERS
# ============================================================================

def _add_title_slide(prs, title, subtitle, palette):
    """Add a dark title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _hex_to_rgb(palette["titleBg"])

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(1.6), Inches(SLIDE_W - 2 * MARGIN), Inches(1.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = palette["headerFont"]
    p.font.color.rgb = _hex_to_rgb(palette["titleText"])
    p.alignment = PP_ALIGN.CENTER

    # Accent line
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(2.85), Inches(3), Inches(0.03),
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = _hex_to_rgb(palette["accent"])
    slide.shapes[-1].line.fill.background()

    # Subtitle
    txBox2 = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(3.0), Inches(SLIDE_W - 2 * MARGIN), Inches(0.6)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.name = palette["bodyFont"]
    p2.font.color.rgb = _hex_to_rgb(palette["accent"])
    p2.alignment = PP_ALIGN.CENTER

    # Date
    today = datetime.now().strftime("%B %Y")
    txBox3 = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(4.8), Inches(SLIDE_W - 2 * MARGIN), Inches(0.4)
    )
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = today
    p3.font.size = Pt(11)
    p3.font.name = palette["bodyFont"]
    p3.font.color.rgb = _hex_to_rgb(palette["titleText"])
    p3.alignment = PP_ALIGN.CENTER


def _add_overview_slide(prs, stages, palette, scope_label):
    """Add the hero overview slide with horizontal flow of stage cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    n = len(stages)
    ov = OverviewLayout

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _hex_to_rgb(palette["slideBg"])

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(ov.TITLE_Y),
        Inches(SLIDE_W - 2 * MARGIN), Inches(ov.TITLE_H),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Value Chain Overview \u2014 {scope_label}"
    p.font.size = Pt(ov.TITLE_FONT_SIZE)
    p.font.bold = True
    p.font.name = palette["headerFont"]
    p.font.color.rgb = _hex_to_rgb(palette["stageText"])

    for i, stage in enumerate(stages):
        x, y, w, h = ov.get_card_position(i, n)

        # Card background (rounded rectangle)
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = _hex_to_rgb(palette["cardBg"])
        card_shape.line.color.rgb = _hex_to_rgb(palette["cardBorder"])
        card_shape.line.width = Pt(0.5)

        # Accent top strip
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(0.06),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = _hex_to_rgb(palette["accent"])
        strip.line.fill.background()

        # Stage number badge (circle)
        badge_size = 0.28
        badge_x = x + (w - badge_size) / 2
        badge_y = y + 0.18
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(badge_x), Inches(badge_y),
            Inches(badge_size), Inches(badge_size),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = _hex_to_rgb(palette["accent"])
        badge.line.fill.background()

        # Number text in badge
        num_box = slide.shapes.add_textbox(
            Inches(badge_x), Inches(badge_y),
            Inches(badge_size), Inches(badge_size),
        )
        num_tf = num_box.text_frame
        num_tf.word_wrap = False
        num_p = num_tf.paragraphs[0]
        num_p.text = str(i + 1)
        num_p.font.size = Pt(ov.STAGE_NUMBER_SIZE)
        num_p.font.bold = True
        num_p.font.name = palette["bodyFont"]
        num_p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        num_p.alignment = PP_ALIGN.CENTER
        num_tf.paragraphs[0].space_before = Pt(0)
        num_tf.paragraphs[0].space_after = Pt(0)

        # Stage name
        name_y = badge_y + badge_size + 0.12
        dynamic_name_size = 8 if n > 7 else (9 if n > 5 else ov.NAME_FONT_SIZE)
        name_box = slide.shapes.add_textbox(
            Inches(x + 0.05), Inches(name_y),
            Inches(w - 0.1), Inches(0.50),
        )
        name_tf = name_box.text_frame
        name_tf.word_wrap = True
        name_p = name_tf.paragraphs[0]
        name_p.text = stage["name"]
        name_p.font.size = Pt(dynamic_name_size)
        name_p.font.bold = True
        name_p.font.name = palette["headerFont"]
        name_p.font.color.rgb = _hex_to_rgb(palette["stageText"])
        name_p.alignment = PP_ALIGN.CENTER

        # Description
        desc_y = name_y + 0.50 + ov.DESC_TOP_PAD
        remaining_h = (y + h) - desc_y - 0.1
        dynamic_desc_size = 7 if n > 7 else ov.DESC_FONT_SIZE
        desc_box = slide.shapes.add_textbox(
            Inches(x + 0.08), Inches(desc_y),
            Inches(w - 0.16), Inches(max(remaining_h, 0.5)),
        )
        desc_tf = desc_box.text_frame
        desc_tf.word_wrap = True
        desc_p = desc_tf.paragraphs[0]
        desc_p.text = stage["description"]
        desc_p.font.size = Pt(dynamic_desc_size)
        desc_p.font.name = palette["bodyFont"]
        desc_p.font.color.rgb = _hex_to_rgb(palette["bodyText"])
        desc_p.alignment = PP_ALIGN.CENTER

        # Arrow to next stage
        if i < n - 1:
            ax, ay, aw, ah = ov.get_arrow_position(i, n)
            arrow_box = slide.shapes.add_textbox(
                Inches(ax), Inches(ay), Inches(aw), Inches(ah),
            )
            arrow_tf = arrow_box.text_frame
            arrow_tf.word_wrap = False
            arrow_p = arrow_tf.paragraphs[0]
            arrow_p.text = "\u203A"  # Single right-pointing angle quotation mark
            arrow_p.font.size = Pt(18)
            arrow_p.font.bold = True
            arrow_p.font.name = palette["bodyFont"]
            arrow_p.font.color.rgb = _hex_to_rgb(palette["arrowColor"])
            arrow_p.alignment = PP_ALIGN.CENTER


def _add_deep_dive_slide(prs, stage, index, total, palette, all_stage_names):
    """Add a deep-dive slide for a single stage."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dd = DeepDiveLayout

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _hex_to_rgb(palette["slideBg"])

    # Stage number badge (large circle)
    badge_size = 0.55
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.5), Inches(dd.STAGE_NUMBER_Y),
        Inches(badge_size), Inches(badge_size),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = _hex_to_rgb(palette["accent"])
    badge.line.fill.background()

    # Number in badge
    num_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(dd.STAGE_NUMBER_Y),
        Inches(badge_size), Inches(badge_size),
    )
    num_tf = num_box.text_frame
    num_p = num_tf.paragraphs[0]
    num_p.text = str(index + 1)
    num_p.font.size = Pt(20)
    num_p.font.bold = True
    num_p.font.name = palette["bodyFont"]
    num_p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    num_p.alignment = PP_ALIGN.CENTER

    # Stage counter (top right)
    counter_box = slide.shapes.add_textbox(
        Inches(SLIDE_W - MARGIN - 1), Inches(dd.STAGE_NUMBER_Y),
        Inches(1), Inches(0.3),
    )
    counter_tf = counter_box.text_frame
    counter_p = counter_tf.paragraphs[0]
    counter_p.text = f"{index + 1} / {total}"
    counter_p.font.size = Pt(dd.STAGE_NUMBER_SIZE)
    counter_p.font.name = palette["bodyFont"]
    counter_p.font.color.rgb = _hex_to_rgb(palette["accent"])
    counter_p.alignment = PP_ALIGN.RIGHT

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(dd.TITLE_X), Inches(dd.TITLE_Y),
        Inches(SLIDE_W - dd.TITLE_X - MARGIN), Inches(dd.TITLE_H),
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = stage["name"]
    title_p.font.size = Pt(dd.TITLE_FONT_SIZE)
    title_p.font.bold = True
    title_p.font.name = palette["headerFont"]
    title_p.font.color.rgb = _hex_to_rgb(palette["stageText"])

    # Description
    desc_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(dd.DESC_Y),
        Inches(SLIDE_W - 2 * MARGIN), Inches(dd.DESC_H),
    )
    desc_tf = desc_box.text_frame
    desc_tf.word_wrap = True
    desc_p = desc_tf.paragraphs[0]
    desc_p.text = stage["description"]
    desc_p.font.size = Pt(dd.DESC_FONT_SIZE)
    desc_p.font.name = palette["bodyFont"]
    desc_p.font.color.rgb = _hex_to_rgb(palette["bodyText"])

    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(MARGIN), Inches(dd.ACCENT_BAR_Y),
        Inches(dd.ACCENT_BAR_W), Inches(dd.ACCENT_BAR_H),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(palette["accent"])
    bar.line.fill.background()

    # "Key Players" label
    label_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(dd.PLAYERS_LABEL_Y),
        Inches(3), Inches(0.35),
    )
    label_tf = label_box.text_frame
    label_p = label_tf.paragraphs[0]
    label_p.text = "Key Players"
    label_p.font.size = Pt(dd.PLAYERS_LABEL_FONT_SIZE)
    label_p.font.bold = True
    label_p.font.name = palette["headerFont"]
    label_p.font.color.rgb = _hex_to_rgb(palette["stageText"])

    # Player cards
    players = stage.get("players", [])
    for p_idx, player in enumerate(players):
        px, py, pw, ph = dd.get_player_card_position(p_idx, len(players))

        # Card bg
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(px), Inches(py), Inches(pw), Inches(ph),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _hex_to_rgb(palette["cardBg"])
        card.line.color.rgb = _hex_to_rgb(palette["cardBorder"])
        card.line.width = Pt(0.5)

        # Left accent bar on card
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(px), Inches(py), Inches(0.06), Inches(ph),
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = _hex_to_rgb(palette["accent"])
        left_bar.line.fill.background()

        # Company name
        name_box = slide.shapes.add_textbox(
            Inches(px + 0.18), Inches(py + 0.1),
            Inches(pw - 0.3), Inches(0.35),
        )
        name_tf = name_box.text_frame
        name_p = name_tf.paragraphs[0]
        name_p.text = player["name"]
        name_p.font.size = Pt(13)
        name_p.font.bold = True
        name_p.font.name = palette["headerFont"]
        name_p.font.color.rgb = _hex_to_rgb(palette["stageText"])

        # Company role
        if player.get("role"):
            role_box = slide.shapes.add_textbox(
                Inches(px + 0.18), Inches(py + 0.45),
                Inches(pw - 0.3), Inches(0.55),
            )
            role_tf = role_box.text_frame
            role_tf.word_wrap = True
            role_p = role_tf.paragraphs[0]
            role_p.text = player["role"]
            role_p.font.size = Pt(9)
            role_p.font.name = palette["bodyFont"]
            role_p.font.color.rgb = _hex_to_rgb(palette["bodyText"])

        # Multi-stage badge
        also_in = player.get("alsoIn", [])
        if also_in:
            also_box = slide.shapes.add_textbox(
                Inches(px + 0.18), Inches(py + ph - 0.25),
                Inches(pw - 0.3), Inches(0.2),
            )
            also_tf = also_box.text_frame
            also_p = also_tf.paragraphs[0]
            also_p.text = f"Also in: {', '.join(also_in)}"
            also_p.font.size = Pt(7)
            also_p.font.italic = True
            also_p.font.name = palette["bodyFont"]
            also_p.font.color.rgb = _hex_to_rgb(palette["accent"])

    # Mini chain breadcrumb at bottom
    mini_y = 4.9
    mini_h = 0.35
    mini_margin = 0.5
    mini_avail_w = SLIDE_W - 2 * mini_margin
    mini_arrow_w = 0.12
    n = len(all_stage_names)
    mini_total_arrows = (n - 1) * mini_arrow_w
    mini_pill_w = (mini_avail_w - mini_total_arrows) / n

    for s in range(n):
        mx = mini_margin + s * (mini_pill_w + mini_arrow_w)
        is_active = s == index

        # Pill background
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(mx), Inches(mini_y),
            Inches(mini_pill_w), Inches(mini_h),
        )
        if is_active:
            pill.fill.solid()
            pill.fill.fore_color.rgb = _hex_to_rgb(palette["accent"])
        else:
            pill.fill.solid()
            pill.fill.fore_color.rgb = _hex_to_rgb(palette["cardBg"])
        pill.line.color.rgb = _hex_to_rgb(palette["accent"] if is_active else palette["cardBorder"])
        pill.line.width = Pt(0.5)

        # Stage name (abbreviated)
        sname = all_stage_names[s]
        short_name = sname[:12] + "\u2026" if len(sname) > 14 else sname
        pill_text = slide.shapes.add_textbox(
            Inches(mx), Inches(mini_y),
            Inches(mini_pill_w), Inches(mini_h),
        )
        pill_tf = pill_text.text_frame
        pill_tf.word_wrap = False
        pill_p = pill_tf.paragraphs[0]
        pill_p.text = short_name
        pill_p.font.size = Pt(6)
        pill_p.font.bold = is_active
        pill_p.font.name = palette["bodyFont"]
        pill_p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if is_active else _hex_to_rgb(palette["bodyText"])
        pill_p.alignment = PP_ALIGN.CENTER

        # Arrow between pills
        if s < n - 1:
            arr_box = slide.shapes.add_textbox(
                Inches(mx + mini_pill_w), Inches(mini_y),
                Inches(mini_arrow_w), Inches(mini_h),
            )
            arr_tf = arr_box.text_frame
            arr_p = arr_tf.paragraphs[0]
            arr_p.text = "\u203A"
            arr_p.font.size = Pt(10)
            arr_p.font.name = palette["bodyFont"]
            arr_p.font.color.rgb = _hex_to_rgb(palette["arrowColor"])
            arr_p.alignment = PP_ALIGN.CENTER


def _add_summary_slide(prs, stages, palette):
    """Add cross-stage players summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sm = SummaryLayout

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _hex_to_rgb(palette["titleBg"])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(sm.TITLE_Y),
        Inches(SLIDE_W - 2 * MARGIN), Inches(sm.TITLE_H),
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "Cross-Stage Players"
    title_p.font.size = Pt(sm.TITLE_FONT_SIZE)
    title_p.font.bold = True
    title_p.font.name = palette["headerFont"]
    title_p.font.color.rgb = _hex_to_rgb(palette["titleText"])

    # Find cross-stage players
    player_stages = {}
    for stage in stages:
        for player in stage.get("players", []):
            name = player["name"]
            if name not in player_stages:
                player_stages[name] = []
            player_stages[name].append(stage["name"])

    cross_stage = [
        {"name": name, "stages": stgs}
        for name, stgs in player_stages.items()
        if len(stgs) > 1
    ]
    cross_stage.sort(key=lambda x: len(x["stages"]), reverse=True)

    if not cross_stage:
        # No cross-stage players, add a note
        note_box = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(sm.TABLE_TOP),
            Inches(SLIDE_W - 2 * MARGIN), Inches(1),
        )
        note_tf = note_box.text_frame
        note_p = note_tf.paragraphs[0]
        note_p.text = "Each company in this value chain operates in a single stage."
        note_p.font.size = Pt(14)
        note_p.font.name = palette["bodyFont"]
        note_p.font.color.rgb = _hex_to_rgb(palette["titleText"])
        note_p.alignment = PP_ALIGN.CENTER
        return

    # Build a visual table using shapes (python-pptx tables look plain)
    row_h = sm.ROW_H
    start_y = sm.TABLE_TOP
    col1_w = sm.COMPANY_COL_W
    col2_w = sm.TABLE_W - col1_w

    # Header row
    for col_x, col_w, text in [
        (sm.TABLE_X, col1_w, "Company"),
        (sm.TABLE_X + col1_w, col2_w, "Stages"),
    ]:
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(col_x), Inches(start_y), Inches(col_w), Inches(row_h),
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = _hex_to_rgb(palette["accent"])
        hdr.line.fill.background()

        hdr_text = slide.shapes.add_textbox(
            Inches(col_x + 0.15), Inches(start_y),
            Inches(col_w - 0.3), Inches(row_h),
        )
        hdr_tf = hdr_text.text_frame
        hdr_p = hdr_tf.paragraphs[0]
        hdr_p.text = text
        hdr_p.font.size = Pt(12)
        hdr_p.font.bold = True
        hdr_p.font.name = palette["headerFont"]
        hdr_p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        hdr_p.alignment = PP_ALIGN.LEFT

    # Data rows (limit to avoid overflow)
    max_rows = min(len(cross_stage), 6)
    for r_idx in range(max_rows):
        cp = cross_stage[r_idx]
        ry = start_y + (r_idx + 1) * row_h
        bg_hex = "1F2937" if r_idx % 2 == 0 else "263040"

        for col_x, col_w, text, is_bold, font_size, color in [
            (sm.TABLE_X, col1_w, cp["name"], True, 11, "FFFFFF"),
            (sm.TABLE_X + col1_w, col2_w, "  \u2192  ".join(cp["stages"]), False, 10, "D1D5DB"),
        ]:
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(col_x), Inches(ry), Inches(col_w), Inches(row_h),
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex_to_rgb(bg_hex)
            cell.line.color.rgb = _hex_to_rgb("374151")
            cell.line.width = Pt(0.5)

            cell_text = slide.shapes.add_textbox(
                Inches(col_x + 0.15), Inches(ry),
                Inches(col_w - 0.3), Inches(row_h),
            )
            cell_tf = cell_text.text_frame
            cell_tf.word_wrap = True
            cell_p = cell_tf.paragraphs[0]
            cell_p.text = text
            cell_p.font.size = Pt(font_size)
            cell_p.font.bold = is_bold
            cell_p.font.name = palette["bodyFont"]
            cell_p.font.color.rgb = _hex_to_rgb(color)


# ============================================================================
# MAIN GENERATOR
# ============================================================================

def build_value_chain_ppt(value_chain, scope="broad"):
    """Build a value chain PowerPoint presentation.

    Args:
        value_chain: dict with 'name', 'broad' and/or 'narrow' keys
        scope: 'broad', 'narrow', or 'both'

    Returns:
        (filepath, filename) tuple
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    chain_name = value_chain["name"]
    palette = _get_palette(chain_name)

    if scope == "both":
        scope_label = "Broad & Narrow"
    elif scope == "narrow":
        focus = value_chain.get("narrow", {}).get("focus", "")
        scope_label = f"Narrow \u2014 {focus}" if focus else "Narrow"
    else:
        scope_label = "Broad"

    # 1. Title slide
    _add_title_slide(prs, f"{chain_name}\nValue Chain", f"{scope_label} Analysis", palette)

    # 2. Build slides for each requested scope
    scopes_to_build = []
    if scope in ("broad", "both"):
        broad = value_chain.get("broad")
        if broad:
            scopes_to_build.append(("Broad", broad["stages"]))
    if scope in ("narrow", "both"):
        narrow = value_chain.get("narrow")
        if narrow:
            label = narrow.get("focus", "Narrow")
            scopes_to_build.append((f"Narrow \u2014 {label}", narrow["stages"]))

    for scope_name, stages in scopes_to_build:
        # Overview slide
        _add_overview_slide(prs, stages, palette, scope_name)

        # Deep-dive slides
        all_names = [s["name"] for s in stages]
        for i, stage in enumerate(stages):
            _add_deep_dive_slide(prs, stage, i, len(stages), palette, all_names)

    # 3. Summary slide (using the last/most detailed set of stages)
    if scopes_to_build:
        last_stages = scopes_to_build[-1][1]
        _add_summary_slide(prs, last_stages, palette)

    # Save
    safe_name = chain_name.replace("/", "-").replace("&", "and").replace(" ", "_")
    filename = f"{safe_name}_Value_Chain_{scope.title()}.pptx"
    fd, filepath = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(filepath)

    return filepath, filename
