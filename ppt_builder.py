"""Build industry landscape PowerPoint presentations."""

import io
import os
import tempfile
import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Logo fetch settings
LOGO_CACHE = {}
LOGO_TIMEOUT = 5  # seconds


def _fetch_logo(domain):
    """Fetch company logo via Clearbit. Returns image bytes or None."""
    if domain in LOGO_CACHE:
        return LOGO_CACHE[domain]
    url = f"https://logo.clearbit.com/{domain}?size=128"
    try:
        resp = requests.get(url, timeout=LOGO_TIMEOUT)
        if resp.status_code == 200 and resp.headers.get("content-type", "").startswith("image"):
            LOGO_CACHE[domain] = resp.content
            return resp.content
    except Exception:
        pass
    LOGO_CACHE[domain] = None
    return None


def _add_title_slide(prs, industry_name, sub_industry_name=None):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)

    # Industry name
    left = Inches(0.8)
    top = Inches(2.0)
    width = Inches(8.4)
    height = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = industry_name
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Sub-industry name
    if sub_industry_name:
        top2 = Inches(3.3)
        height2 = Inches(0.8)
        txBox2 = slide.shapes.add_textbox(left, top2, width, height2)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = sub_industry_name
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
        p2.alignment = PP_ALIGN.CENTER

    return slide


def _add_logo_splash(prs, sub_industry_name, companies):
    """Add a logo splash slide with company logos in a grid."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # White background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Sub-industry header
    header_left = Inches(0.5)
    header_top = Inches(0.3)
    header_w = Inches(9.0)
    header_h = Inches(0.6)
    txBox = slide.shapes.add_textbox(header_left, header_top, header_w, header_h)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = sub_industry_name
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x11, 0x18, 0x27)

    # Grid layout for logos
    n = len(companies)
    if n <= 5:
        cols = n
        rows = 1
    elif n <= 10:
        cols = 5
        rows = 2
    else:
        cols = 5
        rows = 3

    cell_w = Inches(1.6)
    cell_h = Inches(2.2)
    logo_size = Inches(0.9)
    grid_w = cols * cell_w
    start_x = (Inches(10) - grid_w) // 2  # Center horizontally
    start_y = Inches(1.2)

    for idx, company in enumerate(companies[:cols * rows]):
        row = idx // cols
        col = idx % cols
        x = start_x + col * cell_w
        y = start_y + row * cell_h

        # Center of cell
        cx = x + (cell_w - logo_size) // 2
        cy = y

        # Try to add logo
        logo_bytes = _fetch_logo(company["domain"])
        if logo_bytes:
            stream = io.BytesIO(logo_bytes)
            try:
                slide.shapes.add_picture(stream, cx, cy, logo_size, logo_size)
            except Exception:
                _add_placeholder_circle(slide, cx, cy, logo_size, company["name"])
        else:
            _add_placeholder_circle(slide, cx, cy, logo_size, company["name"])

        # Company name below logo
        name_y = cy + logo_size + Inches(0.08)
        name_h = Inches(0.35)
        name_box = slide.shapes.add_textbox(x, name_y, cell_w, name_h)
        ntf = name_box.text_frame
        ntf.word_wrap = True
        np = ntf.paragraphs[0]
        np.text = company["name"]
        np.font.size = Pt(9)
        np.font.bold = True
        np.font.color.rgb = RGBColor(0x11, 0x18, 0x27)
        np.alignment = PP_ALIGN.CENTER

        # Funding below name
        fund_y = name_y + Inches(0.28)
        fund_h = Inches(0.25)
        fund_box = slide.shapes.add_textbox(x, fund_y, cell_w, fund_h)
        ftf = fund_box.text_frame
        ftf.word_wrap = True
        fp = ftf.paragraphs[0]
        fp.text = company.get("funding", "")
        fp.font.size = Pt(7)
        fp.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
        fp.alignment = PP_ALIGN.CENTER

    return slide


def _add_placeholder_circle(slide, x, y, size, name):
    """Add a gray circle with the company's first initial when logo is unavailable."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = name[0].upper() if name else "?"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


def build_landscape_ppt(industry, selected_sub_ids=None):
    """Build a landscape PPT for an industry and selected sub-industries.

    Args:
        industry: dict with 'name' and 'sub_industries' list
        selected_sub_ids: set of sub-industry IDs to include (None = all)

    Returns:
        (filepath, filename) tuple
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # Widescreen 16:9

    industry_name = industry["name"]

    # Title slide for the industry
    _add_title_slide(prs, industry_name)

    # One splash slide per selected sub-industry
    for sub in industry["sub_industries"]:
        if selected_sub_ids and sub["id"] not in selected_sub_ids:
            continue
        companies = sub.get("companies", [])
        if companies:
            _add_logo_splash(prs, sub["name"], companies)

    # Save to temp file
    safe_name = industry_name.replace("/", "-").replace("&", "and").replace(" ", "_")
    filename = f"{safe_name}_Landscape.pptx"
    fd, filepath = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(filepath)

    return filepath, filename
